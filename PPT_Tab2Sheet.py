import os
import sys
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook

#conn = cx_Oracle.connect("DEMAPPRO", "DEMAPPRO", "10.54.58.75:1865/DEM1P")
#cursor = conn.cursor()

BO_Sheet_file = 'BillAnalyzerTestsBO.xlsx'

def main():
    #Open_DB_Connection()
    Process()
    #Close_DB_Connection()

#def Open_DB_Connection():
#    print(conn.version)
#    sqltext = '''select to_char(sysdate,'YYYYMMDDhh24miss') now_text from dual '''
#    cursor.execute(sqltext)
#    
#    for now_text in cursor:
#        print("Start DB Connection: ", now_text)
#
#def Close_DB_Connection():
#    sqltext = '''select to_char(sysdate,'YYYYMMDDhh24miss') now_text from dual '''
#    cursor.execute(sqltext)    
#    for now_text in cursor:
#        print("End DB connection: ", now_text)
#    conn.close()	

def append_df_to_excel(filename, df, sheet_name='Sheet1', startrow=None,            
                       truncate_sheet=False,                                        
                       **to_excel_kwargs):                                          
    """                                                                             
    Append a DataFrame [df] to existing Excel file [filename]                       
    into [sheet_name] Sheet.                                                        
    If [filename] doesn't exist, then this function will create it.                 
                                                                                    
    @param filename: File path or existing ExcelWriter                              
                     (Example: '/path/to/file.xlsx')                                
    @param df: DataFrame to save to workbook                                        
    @param sheet_name: Name of sheet which will contain DataFrame.                  
                       (default: 'Sheet1')                                          
    @param startrow: upper left cell row to dump data frame.                        
                     Per default (startrow=None) calculate the last row             
                     in the existing DF and write to the next row...                
    @param truncate_sheet: truncate (remove and recreate) [sheet_name]              
                           before writing DataFrame to Excel file                   
    @param to_excel_kwargs: arguments which will be passed to `DataFrame.to_excel()`
                            [can be a dictionary]                                   
    @return: None                                                                   
                                                                                    
    Usage examples:                                                                 
                                                                                    
    >>> append_df_to_excel('d:/temp/test.xlsx', df)                                 
                                                                                    
    >>> append_df_to_excel('d:/temp/test.xlsx', df, header=None, index=False)       
                                                                                    
    >>> append_df_to_excel('d:/temp/test.xlsx', df, sheet_name='Sheet2',            
                           index=False)                                             
                                                                                    
    >>> append_df_to_excel('d:/temp/test.xlsx', df, sheet_name='Sheet2',            
                           index=False, startrow=25)                                
                                                                                    
    (c) [MaxU](https://stackoverflow.com/users/5741205/maxu?tab=profile)            
    """                                                                             
    # Excel file doesn't exist - saving and exiting                                 
    if not os.path.isfile(filename):                                                
        df.to_excel(                                                                
            filename,                                                               
            sheet_name=sheet_name,                                                  
            startrow=startrow if startrow is not None else 0,                       
            **to_excel_kwargs)                                                      
        return                                                                      
                                                                                    
    # ignore [engine] parameter if it was passed                                    
    if 'engine' in to_excel_kwargs:                                                 
        to_excel_kwargs.pop('engine')                                               
                                                                                    
    writer = pd.ExcelWriter(filename, engine='openpyxl', mode='a')                  
                                                                                    
    # try to open an existing workbook                                              
    writer.book = load_workbook(filename)                                           
                                                                                    
    # get the last row in the existing Excel sheet                                  
    # if it was not specified explicitly                                            
    if startrow is None and sheet_name in writer.book.sheetnames:                   
        startrow = writer.book[sheet_name].max_row                                  
                                                                                    
    # truncate sheet                                                                
    if truncate_sheet and sheet_name in writer.book.sheetnames:                     
        # index of [sheet_name] sheet                                               
        idx = writer.book.sheetnames.index(sheet_name)                              
        # remove [sheet_name]                                                       
        writer.book.remove(writer.book.worksheets[idx])                             
        # create an empty sheet [sheet_name] using old index                        
        writer.book.create_sheet(sheet_name, idx)                                   
                                                                                    
    # copy existing sheets                                                          
    writer.sheets = {ws.title:ws for ws in writer.book.worksheets}                  
                                                                                    
    if startrow is None:                                                            
        startrow = 0                                                                
                                                                                    
    # write out the new sheet                                                       
    df.to_excel(writer, sheet_name, startrow=startrow, **to_excel_kwargs)           
                                                                                    
    # save the workbook                                                             
    writer.save()                                                                   

def Process():
    in_files = [f for f in os.listdir(".") if (f.endswith('.pptx') and f.startswith('R2'))]
    for in_file in in_files:
        print('****************************')
        print(in_file)
        print('****************************')
        
        test_cycle = in_file[3:(in_file.find("_Migration_"))]
        print("test_cycle: "+test_cycle)
        
        prs = Presentation(in_file)
        
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        #text_runs = []
        BOData_dict={}
        BOName = ''
        BOText = ''
        slide_no = 0
        BOData_dict["FVT"]= test_cycle
        
        for slide in prs.slides:
            for shape in slide.shapes:
                slide_no+=1
                print("Slide No: ", slide_no)     
                if not shape.has_table:
                    continue    
               
                tbl = shape.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                
                for r in range(0, row_count):
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        text_runs = []
                        paragraphs = cell.text_frame.paragraphs 
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
                                #print("00Slide No: ", slide_no, " Row: ", r, " Col: ", c, " text: ",run.text) 
                        print("01Slide No: ", slide_no, " Row: ", r, " Col: ", c, " text: ",''.join(text_runs)) 
                        if (slide_no == 4 and r >0 and r<=6):
                            if (c == 1):
                                BOName = ''.join(text_runs)
                            elif(c==3):
                                BOText = ''.join(text_runs)
                        #elif(slide_no == 8 and r >0 and r<=3):
                        elif(slide_no == 7 and r >0 and r<=3):
                            if (c == 1):
                                BOName = ''.join(text_runs)
                            elif(c==3):
                                BOText = ''.join(text_runs)
                        elif(slide_no == 11 and r >0 and r<=3):
                            if (c == 1):
                                BOName = ''.join(text_runs)
                            elif(c==2):
                                BOText = ''.join(text_runs)
                        #elif(slide_no == 15 and r >0 and r<=3):
                        elif(slide_no == 14 and r >0 and r<=3):
                            if (c == 1):
                                BOName = ''.join(text_runs)
                            elif(c==3):
                                BOText = ''.join(text_runs)
                        elif(slide_no == 18 and r >0 and r<=3):
                            if (c == 1):
                                BOName = ''.join(text_runs)
                            elif(c==2):
                                BOText = ''.join(text_runs)
                                
                    if (BOName != ''):                 
                        BOData_dict[BOName]= BOText   
                                        
        #print(BOData_dict)
        df = pd.DataFrame.from_dict([BOData_dict]) 
        print("Update for file: ", in_file) 
        print(df) 
        append_df_to_excel(BO_Sheet_file, df)  

        # sqltext = 'COMMIT'
        # cursor.execute(sqltext)
        # print('COMMIT for the SI')
        # f.close() 
        os.rename(in_file,in_file+'.done')
        
    
    #sqltext = '''select to_char(sysdate,'YYYYMMDDhh24miss') now_text from dual '''
    #cursor.execute(sqltext)    
    #for now_text in cursor:
    #    print("End of Process Now is:", now_text)


if __name__ == '__main__':
    sys.exit(main())